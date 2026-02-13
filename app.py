import streamlit as st
import litellm
import urllib.request
from pypdf import PdfReader
import auth
from gtts import gTTS
import io
import graphviz
import json
from pptx import Presentation
import extra_streamlit_components as stx

st.set_page_config(page_title="Chatbot - Powered by Open Source LLM")

# Cookie Manager
def get_manager():
    return stx.CookieManager()

cookie_manager = get_manager()

# Application Logic
OLLAMA_BASE_URL = "http://localhost:11434"
MODEL_NAME = "ollama/minimax-m2:cloud"



# -------- PDF TEXT EXTRACTION --------
# -------- PDF TEXT EXTRACTION --------
def extract_pdf_text(uploaded_files):
    text = ""
    for uploaded_file in uploaded_files:
        try:
            reader = PdfReader(uploaded_file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        except Exception as e:
            st.error(f"Error reading {uploaded_file.name}: {e}")
    return text


# -------- LLM RESPONSE --------
def generate_response(messages):
    full_response = ""
    
    # Strip 'ollama/' prefix if present
    api_model = MODEL_NAME.replace("ollama/", "")
    
    url = f"{OLLAMA_BASE_URL}/api/chat"
    headers = {"Content-Type": "application/json"}
    data = {
        "model": api_model,
        "messages": messages,
        "stream": True
    }
    
    try:
        req = urllib.request.Request(url, data=json.dumps(data).encode('utf-8'), headers=headers, method="POST")
        with urllib.request.urlopen(req) as response:
            for line in response:
                if line:
                    decoded_line = line.decode('utf-8').strip()
                    if not decoded_line:
                        continue
                    try:
                        json_obj = json.loads(decoded_line)
                        if "message" in json_obj:
                            content = json_obj["message"].get("content", "")
                            if content:
                                full_response += content
                        # Ignore 'thinking' field or empty content
                    except json.JSONDecodeError:
                        continue
    except Exception as e:
        st.error(f"Error generating response: {e}")
        return f"Error: {e}"

    return full_response


# -------- AUTH & SESSION --------
auth.init_db()

if "user" not in st.session_state:
    st.session_state.user = None

if "messages" not in st.session_state:
    st.session_state.messages = []

if "pdf_context" not in st.session_state:
    st.session_state.pdf_context = ""


def main_app():
    st.sidebar.title(f"Welcome, {st.session_state.user}")
    if st.sidebar.button("Logout"):
        st.session_state.user = None
        cookie_manager.delete("user_token")
        st.rerun()

    # Admin Panel
    if auth.is_admin(st.session_state.user):
        with st.sidebar.expander("Admin Panel"):
            st.write("### User Management")
            users = auth.get_all_users()
            for u in users:
                is_admin_text = " (Admin)" if u[1] else ""
                st.write(f"- {u[0]}{is_admin_text}")

    st.title("💬 Chatbot")
    st.caption("🚀 Streamlit chatbot powered by Ollama (Qwen3-VL 4B)")

    uploaded_pdfs = st.sidebar.file_uploader("Upload PDFs", type=["pdf"], accept_multiple_files=True)
    if uploaded_pdfs:
        pdf_text = extract_pdf_text(uploaded_pdfs)
        st.session_state.pdf_context = pdf_text
        st.sidebar.success(f"{len(uploaded_pdfs)} PDF(s) processed!")

    # Export Chat
    chat_str = "\n".join([f"{m['role']}: {m['content']}" for m in st.session_state.messages])
    st.sidebar.download_button("Download Chat", chat_str, file_name="chat_history.txt")

    # Interactive Features
    with st.sidebar.expander("Interactive Features"):
        feature_mode = st.radio("Select Mode", ["Chat", "Audio Summary", "Quiz", "Knowledge Graph", "Slide Deck"])

    if feature_mode == "Chat":
        for msg in st.session_state.messages:
            st.chat_message(msg["role"]).write(msg["content"])

        if prompt := st.chat_input("Ask something..."):
            st.session_state.messages.append({"role": "user", "content": prompt})
            st.chat_message("user").write(prompt)

            auth.save_message(st.session_state.user, "user", prompt)

            messages_for_llm = st.session_state.messages.copy()
            if st.session_state.pdf_context:
                messages_for_llm.insert(0, {"role": "system", "content": "You are an assistant answering questions based on the following PDF content:\n\n" + st.session_state.pdf_context})

            response = generate_response(messages_for_llm)
            st.session_state.messages.append({"role": "assistant", "content": response})
            st.chat_message("assistant").write(response)
            auth.save_message(st.session_state.user, "assistant", response)

    elif feature_mode == "Audio Summary":
        st.subheader("🔊 Audio Summary")
        if st.button("Generate Audio Summary"):
            if st.session_state.pdf_context:
                with st.spinner("Generating summary and audio..."):
                     # 1. Get Summary
                    messages = [{"role": "system", "content": "Summarize the following text in 3 sentences:"},
                                {"role": "user", "content": st.session_state.pdf_context[:5000]}] # Limit context for speed
                    summary = generate_response(messages)
                    st.write(summary)
                    
                    # 2. Convert to Audio
                    tts = gTTS(text=summary, lang='en')
                    audio_fp = io.BytesIO()
                    tts.write_to_fp(audio_fp)
                    st.audio(audio_fp, format='audio/mp3')
            else:
                st.warning("Please upload a PDF first.")

    elif feature_mode == "Quiz":
        st.subheader("📝 PDF Quiz")
        if st.button("Generate Quiz"):
             if st.session_state.pdf_context:
                with st.spinner("Generating quiz..."):
                    prompt = f"""Generate 3 multiple-choice questions based on the text. 
                    Return ONLY raw JSON (no markdown formatting) in this format:
                    [
                        {{"question": "...", "options": ["A", "B", "C"], "answer": "A"}},
                        ...
                    ]
                    Text: {st.session_state.pdf_context[:3000]}"""
                    
                    messages = [{"role": "user", "content": prompt}]
                    response = generate_response(messages)
                    try:
                        # Clean up potentially markdown-wrapped json
                        cleaned_response = response.replace("```json", "").replace("```", "").strip()
                        st.session_state.quiz_data = json.loads(cleaned_response)
                        st.session_state.quiz_score = 0
                    except Exception as e:
                        st.error(f"Failed to parse quiz JSON: {e}")
                        st.write("Raw response:", response)

        if "quiz_data" in st.session_state:
             with st.form("quiz_form"):
                score = 0
                for i, q in enumerate(st.session_state.quiz_data):
                    st.write(f"**Q{i+1}: {q['question']}**")
                    user_answer = st.radio(f"Select answer for Q{i+1}", q['options'], key=f"q_{i}")
                    if user_answer == q['answer']:
                        score += 1
                
                if st.form_submit_button("Submit"):
                    st.success(f"You scored {score}/{len(st.session_state.quiz_data)}")

    elif feature_mode == "Knowledge Graph":
        st.subheader("🕸️ Knowledge Graph")
        if st.button("Generate Graph"):
             if st.session_state.pdf_context:
                with st.spinner("Analyzing text..."):
                    prompt = f"""Extract top 5 relationships from the text.
                    Return ONLY raw JSON (no markdown) in this format:
                    [
                        {{"source": "Entity1", "target": "Entity2", "label": "relationship"}},
                        ...
                    ]
                    Text: {st.session_state.pdf_context[:3000]}"""
                    
                    messages = [{"role": "user", "content": prompt}]
                    response = generate_response(messages)
                    
                    try:
                        cleaned_response = response.replace("```json", "").replace("```", "").strip()
                        rels = json.loads(cleaned_response)
                        
                        graph = graphviz.Digraph()
                        for r in rels:
                            graph.edge(r['source'], r['target'], label=r['label'])
                        
                        st.graphviz_chart(graph)
                    except Exception as e:
                         st.error(f"Failed to generate graph: {e}")
                         st.write("Raw response:", response)

    elif feature_mode == "Slide Deck":
        st.subheader("📽️ Smart Slide Deck")
        if st.button("Generate Slides"):
             if st.session_state.pdf_context:
                with st.spinner("Generating slide content..."):
                    prompt = f"""Generate a 5-slide presentation based on the text.
                    Return ONLY raw JSON (no markdown) in this format:
                    [
                        {{"title": "Slide Title", "bullets": ["Point 1", "Point 2", "Point 3"]}},
                        ...
                    ]
                    Text: {st.session_state.pdf_context[:4000]}"""
                    
                    messages = [{"role": "user", "content": prompt}]
                    response = generate_response(messages)
                    
                    try:
                        cleaned_response = response.replace("```json", "").replace("```", "").strip()
                        slides_data = json.loads(cleaned_response)
                        
                        prs = Presentation()
                        
                        # Title Slide
                        title_slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        subtitle = slide.placeholders[1]
                        title.text = "Document Summary"
                        subtitle.text = "Generated by IntelliDoc AI"
                        
                        # Content Slides
                        bullet_slide_layout = prs.slide_layouts[1]
                        
                        for slide_data in slides_data:
                            slide = prs.slides.add_slide(bullet_slide_layout)
                            shapes = slide.shapes
                            title_shape = shapes.title
                            body_shape = shapes.placeholders[1]
                            
                            title_shape.text = slide_data.get("title", "Updated Slide")
                            
                            tf = body_shape.text_frame
                            for bullet in slide_data.get("bullets", []):
                                p = tf.add_paragraph()
                                p.text = bullet
                                p.level = 0
                        
                        binary_output = io.BytesIO()
                        prs.save(binary_output)
                        st.download_button("Download PPTX", binary_output.getvalue(), "presentation.pptx")
                        
                    except Exception as e:
                         st.error(f"Failed to generate slides: {e}")
                         st.write("Raw response:", response)


def login_page():
    st.markdown("<h1 style='text-align: center; margin-bottom: 2rem;'>👋 Welcome Back</h1>", unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.subheader("Login to your account")
        username = st.text_input("Username", key="login_user")
        password = st.text_input("Password", type="password", key="login_pass")
        
        if st.button("Login", use_container_width=True):
            if auth.authenticate_user(username, password):
                st.session_state.user = username
                cookie_manager.set("user_token", username)
                history = auth.get_chat_history(username)
                if history:
                    st.session_state.messages = history
                else:
                     st.session_state.messages = [{"role": "assistant", "content": "Hi! You can chat with me or upload a PDF 📄"}]
                st.rerun()
            else:
                st.error("Invalid username or password")


def register_page():
    st.markdown("<h1 style='text-align: center; margin-bottom: 2rem;'>🚀 Create Account</h1>", unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        st.subheader("Join IntelliDoc AI")
        username = st.text_input("New Username", key="reg_user")
        password = st.text_input("New Password", type="password", key="reg_pass")
        
        if st.button("Register", use_container_width=True):
            if auth.register_user(username, password):
                st.success("Registration successful! Please login.")
            else:
                st.error("Username already exists")


# Check for existing session
cookies = cookie_manager.get_all()
if not st.session_state.user and "user_token" in cookies:
    username = cookies["user_token"]
    st.session_state.user = username
    history = auth.get_chat_history(username)
    if history:
        st.session_state.messages = history
    else:
        st.session_state.messages = [{"role": "assistant", "content": "Hi! You can chat with me or upload a PDF 📄"}]


if st.session_state.user:
    main_app()
else:
    # Hero Section
    st.markdown("""
    <div style="text-align: center; padding: 2rem 0;">
        <h1 style="font-size: 3rem;">
            IntelliDoc AI
        </h1>
        <p style="font-size: 1.2rem; opacity: 0.8;">
            Your Advanced Document Assistant. Chat, Analyze, Visualize.
        </p>
    </div>
    """, unsafe_allow_html=True)

    tab1, tab2 = st.tabs(["Login", "Register"])
    with tab1:
        login_page()
    with tab2:
        register_page()
